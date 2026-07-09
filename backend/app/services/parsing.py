import asyncio
import importlib.util
import json
import tempfile
from functools import lru_cache
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from app.config import get_settings
from app.services.chunking import ContentBlock
from app.services.storage import document_artifact_dir


settings = get_settings()


async def parse_document(document_id: str, path: Path, filename: str, mime_type: str = "") -> Tuple[List[ContentBlock], int]:
    blocks = await _try_rag_anything(document_id, path, filename)
    if not blocks:
        blocks = _parse_with_lightweight_parsers(document_id, path, filename, mime_type)
    page_count = max([block.page_number for block in blocks], default=0)
    artifact_path = document_artifact_dir(document_id) / "parsed.json"
    artifact_path.write_text(
        json.dumps([_block_to_dict(block) for block in blocks], ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return blocks, page_count


async def _try_rag_anything(document_id: str, path: Path, filename: str) -> List[ContentBlock]:
    if not settings.enable_rag_anything:
        return []
    try:
        from raganything import RAGAnything, RAGAnythingConfig  # type: ignore
    except Exception:
        return []

    try:
        working_dir = document_artifact_dir(document_id) / "raganything"
        output_dir = document_artifact_dir(document_id) / "raganything-output"
        config = RAGAnythingConfig(
            working_dir=str(working_dir),
            parser=settings.rag_anything_parser,
            parse_method="auto",
            enable_image_processing=True,
            enable_table_processing=True,
            enable_equation_processing=True,
        )
        rag = RAGAnything(config=config)

        process = getattr(rag, "process_document_complete", None)
        if process is None:
            return []
        result = process(file_path=str(path), output_dir=str(output_dir))
        if asyncio.iscoroutine(result):
            await result

        blocks = _load_rag_anything_outputs(document_id, output_dir, filename)
        return blocks
    except Exception:
        return []


def _load_rag_anything_outputs(document_id: str, output_dir: Path, filename: str) -> List[ContentBlock]:
    blocks: List[ContentBlock] = []
    if not output_dir.exists():
        return blocks

    for candidate in output_dir.rglob("*.json"):
        try:
            data = json.loads(candidate.read_text(encoding="utf-8"))
        except Exception:
            continue
        items = data if isinstance(data, list) else data.get("chunks") or data.get("content") or []
        if not isinstance(items, list):
            continue
        for item in items:
            if not isinstance(item, dict):
                continue
            content = str(item.get("content") or item.get("text") or "").strip()
            if not content:
                continue
            blocks.append(
                ContentBlock(
                    document_id=document_id,
                    page_number=int(item.get("page_number") or item.get("page") or 1),
                    block_type=str(item.get("type") or item.get("content_type") or "text"),
                    content=content,
                    metadata={"parser": "rag-anything", "source_file": filename},
                )
            )
    return blocks


def _parse_with_lightweight_parsers(
    document_id: str, path: Path, filename: str, mime_type: str = ""
) -> List[ContentBlock]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(document_id, path, filename)
    if suffix in {".txt", ".md", ".csv"}:
        return [_text_block(document_id, 1, path.read_text(encoding="utf-8", errors="ignore"), filename)]
    if suffix == ".docx":
        return _parse_docx(document_id, path, filename)
    if suffix == ".pptx":
        return _parse_pptx(document_id, path, filename)
    if suffix == ".xlsx":
        return _parse_xlsx(document_id, path, filename)
    if suffix in {".png", ".jpg", ".jpeg"} or mime_type.startswith("image/"):
        return _parse_image(document_id, path, filename)
    return [_text_block(document_id, 1, f"Unsupported file type for fallback parser: {filename}", filename)]


def _parse_pdf(document_id: str, path: Path, filename: str) -> List[ContentBlock]:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        return [_text_block(document_id, 1, f"PDF parser unavailable: {exc}", filename)]

    blocks: List[ContentBlock] = []
    reader = PdfReader(str(path))
    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            blocks.append(_text_block(document_id, page_index, text, filename, parser="pypdf"))
    if not blocks:
        blocks = _ocr_pdf(document_id, path, filename, len(reader.pages))
    return blocks


def _ocr_pdf(document_id: str, path: Path, filename: str, page_count: int) -> List[ContentBlock]:
    if not settings.pdf_ocr_enabled:
        return [_text_block(document_id, 1, "No extractable PDF text found. OCR parser is disabled.", filename)]

    try:
        import pypdfium2 as pdfium
    except Exception as exc:
        return [_text_block(document_id, 1, f"No extractable PDF text found. PDF renderer unavailable: {exc}", filename)]

    blocks: List[ContentBlock] = []
    try:
        pdf = pdfium.PdfDocument(str(path))
        max_pages = min(len(pdf), max(1, settings.pdf_ocr_max_pages))
        for page_index in range(max_pages):
            page = pdf[page_index]
            bitmap = page.render(scale=settings.pdf_ocr_scale)
            image = bitmap.to_pil()
            blocks.append(_ocr_image_to_block(document_id, image, filename, page_index + 1))
    except Exception as exc:
        return [_text_block(document_id, 1, f"No extractable PDF text found. OCR failed: {exc}", filename)]

    if page_count > len(blocks):
        blocks.append(
            _text_block(
                document_id,
                len(blocks) + 1,
                f"OCR stopped after {len(blocks)} of {page_count} pages due to PDF_OCR_MAX_PAGES.",
                filename,
                parser="tesseract-pdf-ocr",
            )
        )
    return blocks or [_text_block(document_id, 1, "No extractable PDF text found. OCR produced no pages.", filename)]


def _ocr_image_to_block(document_id: str, image, filename: str, page_number: int) -> ContentBlock:
    if settings.pdf_ocr_engine in {"auto", "paddle"}:
        paddle_result = _ocr_image_with_paddle(image)
        if paddle_result is not None and (paddle_result["text"].strip() or settings.pdf_ocr_engine == "paddle"):
            return ContentBlock(
                document_id=document_id,
                page_number=page_number,
                block_type="text",
                content=paddle_result["text"].strip() or "OCR produced no text for this page.",
                metadata={
                    "parser": "paddleocr",
                    "source_file": filename,
                    "ocr_lang": settings.paddle_ocr_lang,
                    "ocr_device": settings.paddle_ocr_device,
                    "ocr_scale": settings.pdf_ocr_scale,
                    "ocr_confidence_avg": paddle_result["confidence_avg"],
                    "ocr_confidence_min": paddle_result["confidence_min"],
                    "ocr_line_count": len(paddle_result["lines"]),
                    "ocr_lines": paddle_result["lines"],
                },
            )

    if settings.pdf_ocr_engine == "paddle":
        return _text_block(document_id, page_number, "PaddleOCR unavailable or produced no text.", filename, parser="paddleocr")

    if settings.pdf_ocr_engine in {"auto", "vietocr"}:
        vietocr_result = _ocr_image_with_vietocr(image)
        if vietocr_result is not None and (vietocr_result["text"].strip() or settings.pdf_ocr_engine == "vietocr"):
            return ContentBlock(
                document_id=document_id,
                page_number=page_number,
                block_type="text",
                content=vietocr_result["text"].strip() or "OCR produced no text for this page.",
                metadata={
                    "parser": "vietocr",
                    "source_file": filename,
                    "ocr_device": settings.vietocr_device,
                    "ocr_config": settings.vietocr_config,
                    "ocr_scale": settings.pdf_ocr_scale,
                    "ocr_confidence_avg": vietocr_result["confidence_avg"],
                    "ocr_confidence_min": vietocr_result["confidence_min"],
                    "ocr_line_count": len(vietocr_result["lines"]),
                    "ocr_lines": vietocr_result["lines"],
                },
            )

    if settings.pdf_ocr_engine == "vietocr":
        return _text_block(document_id, page_number, "VietOCR unavailable or produced no text.", filename, parser="vietocr")
    return _ocr_image_with_tesseract(document_id, image, filename, page_number, fallback_reason="advanced_ocr_unavailable_or_empty")


def _ocr_image_with_paddle(image) -> Optional[Dict[str, Any]]:
    try:
        ocr = _paddle_ocr_engine()
    except Exception:
        return None

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_image:
        temp_path = Path(temp_image.name)
    try:
        image.save(temp_path)
        result = _run_paddle_ocr(ocr, temp_path)
        lines = _filter_ocr_lines(_extract_paddle_lines(result))
        text = "\n".join(line["text"] for line in lines if line["text"].strip())
        confidences = [line["confidence"] for line in lines if line.get("confidence") is not None]
        confidence_avg = round(sum(confidences) / len(confidences), 4) if confidences else None
        confidence_min = round(min(confidences), 4) if confidences else None
        return {
            "text": text,
            "confidence_avg": confidence_avg,
            "confidence_min": confidence_min,
            "lines": lines,
        }
    except Exception:
        return None
    finally:
        temp_path.unlink(missing_ok=True)


def _ocr_image_with_vietocr(image) -> Optional[Dict[str, Any]]:
    try:
        detector = _vietocr_engine()
        import pytesseract

        data = pytesseract.image_to_data(image, lang=settings.pdf_ocr_lang, output_type=pytesseract.Output.DICT)
        boxes, base_confidences = _extract_tesseract_boxes(data)
        lines = []
        for box, base_confidence in zip(boxes, base_confidences):
            left, top, right, bottom = box
            crop = image.crop((left, top, right, bottom))
            text = str(detector.predict(crop)).strip()
            if not text:
                continue
            lines.append({"text": text, "confidence": base_confidence, "bbox": box})
        lines = _filter_ocr_lines(lines)
        confidences = [line["confidence"] for line in lines if line.get("confidence") is not None]
        return {
            "text": "\n".join(line["text"] for line in lines),
            "confidence_avg": round(sum(confidences) / len(confidences), 4) if confidences else None,
            "confidence_min": round(min(confidences), 4) if confidences else None,
            "lines": lines[:200],
        }
    except Exception:
        return None


@lru_cache(maxsize=1)
def _vietocr_engine():
    from vietocr.tool.config import Cfg
    from vietocr.tool.predictor import Predictor

    config = Cfg.load_config_from_name(settings.vietocr_config)
    config["device"] = settings.vietocr_device
    config["predictor"]["beamsearch"] = False
    return Predictor(config)


@lru_cache(maxsize=1)
def _paddle_ocr_engine():
    from paddleocr import PaddleOCR

    kwargs = {
        "lang": settings.paddle_ocr_lang,
        "device": settings.paddle_ocr_device,
        "use_doc_orientation_classify": False,
        "use_doc_unwarping": False,
        "use_textline_orientation": False,
    }
    try:
        return PaddleOCR(**kwargs)
    except TypeError:
        return PaddleOCR(lang=settings.paddle_ocr_lang, use_angle_cls=False)


def _run_paddle_ocr(ocr, image_path: Path):
    predict = getattr(ocr, "predict", None)
    if predict is not None:
        return predict(str(image_path))
    return ocr.ocr(str(image_path), cls=False)


def _extract_paddle_lines(result) -> List[Dict[str, Any]]:
    lines: List[Dict[str, Any]] = []
    for page_result in result or []:
        if isinstance(page_result, dict):
            lines.extend(_extract_paddle_lines_from_dict(page_result))
        elif isinstance(page_result, list):
            lines.extend(_extract_paddle_lines_from_legacy(page_result))
        elif hasattr(page_result, "json"):
            try:
                lines.extend(_extract_paddle_lines_from_dict(page_result.json))
            except Exception:
                continue
    return lines


def _extract_paddle_lines_from_dict(result: Dict[str, Any]) -> List[Dict[str, Any]]:
    texts = result.get("rec_texts") or []
    raw_scores = _first_present(result, "rec_scores")
    raw_polygons = _first_present(result, "rec_polys", "dt_polys", "rec_boxes")
    scores = _to_plain_list(raw_scores) if raw_scores is not None else []
    polygons = _to_plain_list(raw_polygons) if raw_polygons is not None else []
    lines = []
    for index, text in enumerate(texts):
        confidence = _safe_float(scores[index]) if index < len(scores) else None
        lines.append(
            {
                "text": str(text),
                "confidence": confidence,
                "bbox": _to_plain_list(polygons[index]) if index < len(polygons) else None,
            }
        )
    return lines


def _extract_paddle_lines_from_legacy(result: List[Any]) -> List[Dict[str, Any]]:
    lines = []
    for item in result:
        if not isinstance(item, list) or len(item) < 2:
            continue
        text_score = item[1]
        if not isinstance(text_score, (list, tuple)) or not text_score:
            continue
        lines.append(
            {
                "text": str(text_score[0]),
                "confidence": _safe_float(text_score[1]) if len(text_score) > 1 else None,
                "bbox": _to_plain_list(item[0]),
            }
        )
    return lines


def _ocr_image_with_tesseract(
    document_id: str,
    image,
    filename: str,
    page_number: int,
    fallback_reason: str = None,
) -> ContentBlock:
    try:
        import pytesseract

        data = pytesseract.image_to_data(image, lang=settings.pdf_ocr_lang, output_type=pytesseract.Output.DICT)
        lines, confidence_avg, confidence_min = _extract_tesseract_lines(data)
        content = "\n".join(line["text"] for line in lines).strip() or "OCR produced no text for this page."
    except Exception as exc:
        content = f"OCR unavailable for this page: {exc}"
        lines, confidence_avg, confidence_min = [], None, None
    return ContentBlock(
        document_id=document_id,
        page_number=page_number,
        block_type="text",
        content=content,
        metadata={
            "parser": "tesseract-pdf-ocr",
            "source_file": filename,
            "ocr_lang": settings.pdf_ocr_lang,
            "ocr_scale": settings.pdf_ocr_scale,
            "ocr_confidence_avg": confidence_avg,
            "ocr_confidence_min": confidence_min,
            "ocr_line_count": len(lines),
            "ocr_lines": lines[:200],
            "preferred_ocr_engine": settings.pdf_ocr_engine,
            "paddleocr_available": _module_available("paddleocr"),
            "vietocr_available": _module_available("vietocr"),
            "fallback_reason": fallback_reason,
        },
    )


def _extract_tesseract_lines(data: Dict[str, List[Any]]) -> Tuple[List[Dict[str, Any]], Optional[float], Optional[float]]:
    grouped: Dict[Tuple[int, int, int], Dict[str, Any]] = {}
    confidences = []
    for index, text in enumerate(data.get("text", [])):
        clean_text = str(text).strip()
        confidence = _safe_float(data.get("conf", [None])[index])
        if not clean_text or confidence is None or confidence < 0:
            continue
        key = (
            int(data.get("block_num", [0])[index]),
            int(data.get("par_num", [0])[index]),
            int(data.get("line_num", [0])[index]),
        )
        left = int(data.get("left", [0])[index])
        top = int(data.get("top", [0])[index])
        width = int(data.get("width", [0])[index])
        height = int(data.get("height", [0])[index])
        line = grouped.setdefault(key, {"words": [], "confidences": [], "bbox": [left, top, left + width, top + height]})
        line["words"].append(clean_text)
        line["confidences"].append(confidence / 100.0)
        line["bbox"][0] = min(line["bbox"][0], left)
        line["bbox"][1] = min(line["bbox"][1], top)
        line["bbox"][2] = max(line["bbox"][2], left + width)
        line["bbox"][3] = max(line["bbox"][3], top + height)
        confidences.append(confidence / 100.0)

    lines = _filter_ocr_lines([
        {
            "text": " ".join(value["words"]),
            "confidence": round(sum(value["confidences"]) / len(value["confidences"]), 4),
            "bbox": value["bbox"],
        }
        for value in grouped.values()
        if value["words"]
    ])
    filtered_confidences = [line["confidence"] for line in lines if line.get("confidence") is not None]
    confidence_avg = round(sum(filtered_confidences) / len(filtered_confidences), 4) if filtered_confidences else None
    confidence_min = round(min(filtered_confidences), 4) if filtered_confidences else None
    return lines[:200], confidence_avg, confidence_min


def _extract_tesseract_boxes(data: Dict[str, List[Any]]) -> Tuple[List[List[int]], List[Optional[float]]]:
    lines, _, _ = _extract_tesseract_lines(data)
    return [line["bbox"] for line in lines], [line.get("confidence") for line in lines]


def _filter_ocr_lines(lines: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    filtered = []
    for line in lines:
        confidence = line.get("confidence")
        if confidence is not None and confidence < settings.ocr_min_line_confidence:
            continue
        text = str(line.get("text") or "").strip()
        if not text:
            continue
        filtered.append({**line, "text": text})
    return filtered


def _to_plain_list(value):
    if hasattr(value, "tolist"):
        return value.tolist()
    if isinstance(value, tuple):
        return [_to_plain_list(item) for item in value]
    if isinstance(value, list):
        return [_to_plain_list(item) for item in value]
    return value


def _first_present(data: Dict[str, Any], *keys: str):
    for key in keys:
        value = data.get(key)
        if value is not None:
            return value
    return None


def _safe_float(value) -> Optional[float]:
    try:
        return round(float(value), 4)
    except (TypeError, ValueError):
        return None


def _module_available(module_name: str) -> bool:
    return importlib.util.find_spec(module_name) is not None


def _parse_docx(document_id: str, path: Path, filename: str) -> List[ContentBlock]:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    table_text = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_text.append(" | ".join(cells))
    content = "\n\n".join(paragraphs + table_text)
    return [_text_block(document_id, 1, content or "Empty DOCX document.", filename, parser="python-docx")]


def _parse_pptx(document_id: str, path: Path, filename: str) -> List[ContentBlock]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    blocks: List[ContentBlock] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        blocks.append(_text_block(document_id, index, "\n".join(texts), filename, parser="python-pptx"))
    return blocks


def _parse_xlsx(document_id: str, path: Path, filename: str) -> List[ContentBlock]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    blocks: List[ContentBlock] = []
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        rows = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 200:
                rows.append("...")
                break
            values = [str(value) for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        blocks.append(
            ContentBlock(
                document_id=document_id,
                page_number=sheet_index,
                block_type="table",
                content=f"Sheet: {sheet.title}\n" + "\n".join(rows),
                metadata={"parser": "openpyxl", "source_file": filename, "sheet": sheet.title},
            )
        )
    return blocks


def _parse_image(document_id: str, path: Path, filename: str) -> List[ContentBlock]:
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(str(path))
        text = pytesseract.image_to_string(image, lang="vie+eng")
        content = text.strip() or f"Image file {filename} produced no OCR text."
    except Exception as exc:
        content = f"Image OCR unavailable for {filename}: {exc}"
    return [
        ContentBlock(
            document_id=document_id,
            page_number=1,
            block_type="image_caption",
            content=content,
            metadata={"parser": "pytesseract", "source_file": filename},
        )
    ]


def _text_block(
    document_id: str,
    page_number: int,
    content: str,
    filename: str,
    parser: str = "fallback",
) -> ContentBlock:
    return ContentBlock(
        document_id=document_id,
        page_number=page_number,
        block_type="text",
        content=content,
        metadata={"parser": parser, "source_file": filename},
    )


def _block_to_dict(block: ContentBlock):
    return {
        "document_id": block.document_id,
        "page_number": block.page_number,
        "block_type": block.block_type,
        "content": block.content,
        "metadata": block.metadata,
    }
